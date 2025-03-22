import pandas as pd
from docx import Document
from bs4 import BeautifulSoup
import markdown2
from odf import opendocument, text
from PyPDF2 import PdfReader
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
from openpyxl import load_workbook

class TextExtractor:
    """
    A class to extract text from various file formats including CSV, DOCX, HTML, Markdown, PDF, PPTX, ODT, RTF, and TXT.
    """
    
    def __init__(self, logger=None):
        self.logger = logger

    
    def extract_text_from_csv(self, file_path):
        """
        Extract text from a CSV file.
        
        Args:
            file_path (str): Path to the CSV file.
        
        Returns:
            str: Text content from the CSV file.
        """
        self.logger.info(f"Extracting text from CSV file: {file_path}")
        df = pd.read_csv(file_path)
        return df.to_string(index=False)

    
    def extract_text_from_docx(self, file_path):
        """
        Extract text from a DOCX file.
        
        Args:
            file_path (str): Path to the DOCX file.
        
        Returns:
            str: Text content from the DOCX file.
        """
        self.logger.info(f"Extracting text from DOCX file: {file_path}")
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    
    def extract_text_from_html(self, file_path):
        """
        Extract text from an HTML file.
        
        Args:
            file_path (str): Path to the HTML file.
        
        Returns:
            str: Text content from the HTML file.
        """
        self.logger.info(f"Extracting text from HTML file: {file_path}")
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
            return soup.get_text()

    
    def extract_text_from_md(self, file_path):
        """
        Extract text from a Markdown file.
        
        Args:
            file_path (str): Path to the Markdown file.
        
        Returns:
            str: Text content from the Markdown file.
        """
        self.logger.info(f"Extracting text from Markdown file: {file_path}")
        with open(file_path, 'r', encoding='utf-8') as file:
            html = markdown2.markdown(file.read())
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text()

    
    def extract_text_from_odt(self, file_path):
        """
        Extract text from an ODT file.
        
        Args:
            file_path (str): Path to the ODT file.
        
        Returns:
            str: Text content from the ODT file.
        """
        self.logger.info(f"Extracting text from ODT file: {file_path}")
        doc = opendocument.load(file_path)
        text_content = []
        for element in doc.getElementsByType(text.P):
            text_content.append(str(element))
        return "\n".join(text_content)

    
    def extract_text_from_pdf(self, file_path):
        """
        Extract text from a PDF file.
        
        Args:
            file_path (str): Path to the PDF file.
        
        Returns:
            str: Text content from the PDF file.
        """
        self.logger.info(f"Extracting text from PDF file: {file_path}")
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        return "\n".join(text)

    
    def extract_text_from_pptx(self, file_path):
        """
        Extract text from a PPTX file.
        
        Args:
            file_path (str): Path to the PPTX file.
        
        Returns:
            str: Text content from the PPTX file.
        """
        self.logger.info(f"Extracting text from PPTX file: {file_path}")
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    
    def extract_text_from_rtf(self, file_path):
        """
        Extract text from an RTF file.
        
        Args:
            file_path (str): Path to the RTF file.
        
        Returns:
            str: Text content from the RTF file.
        """
        self.logger.info(f"Extracting text from RTF file: {file_path}")
        with open(file_path, 'r', encoding='utf-8') as file:
            return rtf_to_text(file.read())

    
    def extract_text_from_txt(self, file_path):
        """
        Extract text from a TXT file.
        
        Args:
            file_path (str): Path to the TXT file.
        
        Returns:
            str: Text content from the TXT file.
        """
        self.logger.info(f"Extracting text from TXT file: {file_path}")
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    
    def extract_text_from_xlsx(self, file_path):
        """
        Extract text from an XLSX file.
        
        Args:
            file_path (str): Path to the XLSX file.
        
        Returns:
            str: Text content from the XLSX file.
        """
        self.logger.info(f"Extracting text from XLSX file: {file_path}")
        workbook = load_workbook(file_path, data_only=True)
        text = []
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                text.extend([str(cell) for cell in row if cell is not None])
        return "\n".join(text)
